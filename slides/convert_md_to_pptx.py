#!/usr/bin/env python3
"""Простой конвертер Markdown слайдов в PPTX.

Ожидает входной файл, разделителем слайдов служит строка с тремя дефисами на отдельной строке (`---`).
Первый блок YAML-фронтмета удаляется.

Usage:
  python3 slides/convert_md_to_pptx.py slides/presentation.md slides/presentation.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def parse_slides(md_text: str):
    # remove initial YAML frontmatter if present
    parts = md_text.split('\n')
    if parts and parts[0].strip() == '---':
        # find next '---'
        try:
            idx = parts[1:].index('---') + 1
            parts = parts[idx+1:]
        except ValueError:
            # no closing frontmatter, keep whole
            parts = parts
    body = '\n'.join(parts)
    slides = [s.strip() for s in body.split('\n---\n') if s.strip()]
    return slides


def slide_title_and_body(slide_md: str):
    lines = [l.rstrip() for l in slide_md.split('\n') if l.strip()]
    title = ''
    body_lines = []
    if lines:
        # first header-like line becomes title
        if lines[0].lstrip().startswith('#'):
            title = lines[0].lstrip().lstrip('#').strip()
            body_lines = lines[1:]
        else:
            # fallback: first line = title
            title = lines[0]
            body_lines = lines[1:]
    return title, '\n'.join(body_lines)


def add_textbox(slide, left, top, width, height, text, font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)


def md_to_pptx(input_md: Path, output_pptx: Path):
    md_text = input_md.read_text(encoding='utf-8')
    slides_md = parse_slides(md_text)

    prs = Presentation()
    # default slide layout index 5 is blank in many templates; fallback to 1
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]

    from pptx.util import Inches

    for s in slides_md:
        title, body = slide_title_and_body(s)
        slide = prs.slides.add_slide(blank_layout)
        if title:
            add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(1), title, font_size=28)
        if body:
            add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5), body, font_size=18)

    prs.save(str(output_pptx))


def main():
    if len(sys.argv) < 3:
        print('Usage: convert_md_to_pptx.py input.md output.pptx')
        sys.exit(2)
    input_md = Path(sys.argv[1])
    output_pptx = Path(sys.argv[2])
    if not input_md.exists():
        print('Input file not found:', input_md)
        sys.exit(1)
    md_to_pptx(input_md, output_pptx)
    print('Saved:', output_pptx)


if __name__ == '__main__':
    main()
